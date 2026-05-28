#!/usr/bin/env python3
"""
通用科室会PPT生成器 — 产品无关版
用法:
  # 交互式输入
  python3 generate_dept_meeting_ppt.py

  # 全参数模式
  python3 generate_dept_meeting_ppt.py \\
    --product "帕博利珠单抗" \\
    --indication "非小细胞肺癌" \\
    --audience 肿瘤内科

  # 自定义数据文件
  python3 generate_dept_meeting_ppt.py \\
    --product "XX" --indication "XX" --audience 肿瘤内科 \\
    --data-file ./my-data.json
"""

import sys
import os
import json
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 设计系统 ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_SLIDE = RGBColor(0xF5, 0xF5, 0xFA)
ACCENT = RGBColor(0x00, 0x96, 0xD6)         # 天蓝主色
ACCENT2 = RGBColor(0xFF, 0x6B, 0x35)         # 橙色强调
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_TEXT = RGBColor(0x66, 0x66, 0x66)
TABLE_HEADER = RGBColor(0x00, 0x7A, 0xBC)
TABLE_ROW1 = RGBColor(0xE8, 0xF4, 0xFD)
TABLE_ROW2 = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_GOOD = RGBColor(0x27, 0xAE, 0x60)
RED_BAD = RGBColor(0xE7, 0x4C, 0x3C)
HIGHLIGHT_BG = RGBColor(0xFF, 0xF3, 0xE0)

# 默认占位数据 (无产品依赖的通用模板)
DEFAULT_DATA = {
    "product_name": "产品A",
    "product_code": "CODE-001",
    "drug_class": "靶向治疗药物",
    "mechanism": "通过特异性靶向肿瘤细胞表面高表达的受体，抑制下游信号通路，诱导肿瘤细胞凋亡",
    "target": "xx受体",
    "payload": "（如ADC载荷）",
    "study_code": "XXX-001",
    "disease": "xx肿瘤",
    "incidence": "约XX/10万",
    "key_data": {
        "ORR": "XX%",
        "mPFS": "XX个月",
        "HR": "0.XX",
        "DCR": "XX%",
        "control_ORR": "XX%",
        "control_mPFS": "XX个月",
        "control_DCR": "XX%"
    },
    "arm_name": "试验组",
    "control_arm_name": "对照组",
    "study_type": "多中心、随机、III期",
    "centers": "XX家",
    "enrollment": "XX例",
    "subgroup_notes": "✅ 各亚组一致获益",
    "safety_highlights": [
        "常见AE为XX级",
        "停药率约X%",
        "整体可管理"
    ],
    "comparator_data": [
        ["药物/方案", "类型", "阶段", "ORR", "mPFS"],
        ["本品", "xx", "III期", "XX%", "XXmo"],
        ["对照方案", "xx", "已有数据", "XX%", "XXmo"]
    ],
    "take_home": [
        "临床数据表现优于现有方案，为患者提供新的治疗选择",
        "安全性可管理，不良反应谱可预测",
        "建议在适当患者中积极考虑使用"
    ]
}

# 受众科室
AUDIENCES = ["肿瘤内科", "放疗科", "呼吸科", "胸外科", "乳腺外科",
             "泌尿外科", "消化科", "妇科肿瘤", "血液科", "综合"]

# 受众定制内容模板
AUDIENCE_TEMPLATES = {
    "肿瘤内科": {
        "title": "对肿瘤内科医生的核心价值",
        "items": [
            "肿瘤内科是化疗/免疫治疗的执行科室",
            "患者后线治疗直接由肿瘤内科管理",
            "本方案提供了后线治疗的新选择",
            "安全性管理策略可参照现有经验",
            "",
            "📊 疗效预期：",
            "  标准治疗基础上进一步改善疗效",
            "  PFS延长、疾病控制率提升",
            "",
            "🏆 后线治疗的新选择"
        ]
    },
    "放疗科": {
        "title": "对放疗科医生的核心价值",
        "items": [
            "放疗科是患者长期随访的主力科室",
            "放疗与系统治疗的时序安排值得关注",
            "联合治疗可能是未来的方向",
            "",
            "⚠️ 不良反应鉴别：",
            "  需要与放疗相关不良反应鉴别",
            "  各科室协同管理患者",
            "",
            "📋 适合患者的筛选标准"
        ]
    },
    "综合": {
        "title": "对临床科室的核心价值",
        "items": [
            "本方案在关键临床研究中显示了显著疗效",
            "疗效数据优于现有标准治疗",
            "安全性可管理，不良反应谱清晰",
            "",
            "🏥 适合患者画像：",
            "  多线治疗后进展的患者",
            "  ECOG 0-1，追求有效治疗方案",
            "",
            "📈 未来展望：",
            "  进一步拓展适应症（进行中）",
            "  联合治疗方案探索",
            "  指南纳入→临床普及"
        ]
    }
}

# 默认科室内容补充
for aud in AUDIENCES:
    if aud not in AUDIENCE_TEMPLATES:
        AUDIENCE_TEMPLATES[aud] = AUDIENCE_TEMPLATES["综合"]

prs = None

def init_presentation():
    global prs
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

def add_bg(slide, color=BG_SLIDE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, color=DARK_TEXT,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=13, color=DARK_TEXT,
                    font_name='微软雅黑', line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        from pptx.oxml.ns import qn
        try:
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
        except:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
    return txBox

def add_kpi_box(slide, left, top, width, height, number, label, num_color=ACCENT):
    add_shape_bg(slide, WHITE, left, top, width, height)
    add_shape_bg(slide, num_color, left, top, width, Pt(4))
    add_textbox(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.5),
                number, font_size=28, color=num_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.05), top + Inches(0.55), width - Inches(0.1), Inches(0.4),
                label, font_size=10, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_data in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_data)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = '微软雅黑'
                paragraph.alignment = PP_ALIGN.CENTER
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.size = Pt(11)
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW1
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW2
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_cover(slide, data, audience):
    bg = add_shape_bg(slide, BG_DARK, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    add_shape_bg(slide, ACCENT, Emu(0), Emu(0), Emu(80000), prs.slide_height)
    add_textbox(slide, Inches(1.0), Inches(1.2), Inches(10), Inches(1.2),
                f'{data["indication"]}治疗新选择',
                font_size=32, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.0), Inches(2.4), Inches(10), Inches(0.8),
                f'{data["product_name"]} ({data.get("product_code", "")}) {data.get("study_code", "")}临床研究解读',
                font_size=20, color=ACCENT)
    add_shape_bg(slide, ACCENT2, Inches(1.0), Inches(3.4), Inches(2.5), Pt(3))
    add_textbox(slide, Inches(1.0), Inches(3.7), Inches(6), Inches(0.5),
                f'产品类型：{data["drug_class"]} · 靶点：{data["target"]}',
                font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_textbox(slide, Inches(1.0), Inches(4.3), Inches(8), Inches(1.5),
                f'{data.get("study_type", "多中心临床研究")}\n'
                f'{data.get("centers", "XX家")} | {data.get("enrollment", "XX例")}',
                font_size=13, color=RGBColor(0xCC, 0xCC, 0xCC))
    add_textbox(slide, Inches(1.0), Inches(5.8), Inches(8), Inches(0.5),
                f'学术目标科室：{audience} · 科室内部学术会议',
                font_size=14, color=ACCENT)

def add_section_header(slide, number, title):
    add_bg(slide, BG_SLIDE)
    add_shape_bg(slide, ACCENT, Emu(0), Emu(0), prs.slide_width, Emu(55000))
    add_textbox(slide, Inches(0.5), Inches(0.05), Inches(10), Inches(0.5),
                f'{number}. {title}', font_size=22, color=WHITE, bold=True)

def build_standard_deck(data, audience="综合"):
    kd = data.get("key_data", {})
    audience_content = AUDIENCE_TEMPLATES.get(audience, AUDIENCE_TEMPLATES["综合"])

    # Slide 1: 封面
    add_cover(new_slide(), data, audience)

    # Slide 2: 疾病背景与未满足需求
    s = new_slide()
    add_section_header(s, '1', '未满足的临床需求')
    background_items = [
        f'🌏 全球{data["disease"]}发病率约为{data["incidence"]}',
        '🎯 当前一线标准治疗存在局限性',
        '⚠️ 仍有相当比例患者出现疾病进展',
        '❌ 后线治疗选择有限，亟需更有效方案',
        '',
        '💊 现有治疗现状：',
        '  一线治疗：化疗/靶向/免疫治疗',
        '  后线治疗：可选择方案有限',
        '  疗效瓶颈：ORR/PFS均有提升空间',
        '',
        '⬇ 亟需更有效、更持久的治疗方案'
    ]
    add_bullet_text(s, Inches(0.5), Inches(0.8), Inches(5.5), Inches(4.0), background_items, font_size=12)
    add_shape_bg(s, RGBColor(0xFF, 0xEB, 0xEE), Inches(6.5), Inches(0.8), Inches(5.0), Inches(3.8))
    add_textbox(s, Inches(6.8), Inches(0.9), Inches(4.5), Inches(0.4),
                '未被满足的需求', font_size=16, color=RED_BAD, bold=True)
    need_items = [
        '🔴 后线治疗数据多为单臂/小样本',
        '🔴 后线标准治疗的疗效有限',
        '🔴 多数患者在现有方案进展后选择有限',
        '🔴 缺乏高质量的III期RCT证据',
        '',
        '⬇ 亟需更多循证医学证据'
    ]
    add_bullet_text(s, Inches(6.8), Inches(1.4), Inches(4.5), Inches(2.5), need_items, font_size=12)

    # Slide 3: 药物机制
    s = new_slide()
    add_section_header(s, '2', f'{data["product_name"]}：作用机制')
    add_shape_bg(s, WHITE, Inches(0.5), Inches(0.8), Inches(5.5), Inches(4.5))
    add_textbox(s, Inches(0.7), Inches(0.9), Inches(5.0), Inches(0.4),
                f'{data["drug_class"]} — 作用机制', font_size=14, color=DARK_TEXT, bold=True)
    mech_items = [
        f'🎯 靶点：{data["target"]}',
        f'🧬 抗体制剂：{data.get("antibody_type", "单克隆抗体")}',
        f'{"💊 载荷：" + data["payload"] if data.get("payload") and data["payload"] != "（如ADC载荷）" else ""}',
        '',
        '📐 作用机制：',
        f'  {data["mechanism"]}',
        '',
        '✅ 关键优势：精准靶向，减少脱靶毒性',
    ]
    mech_items = [item for item in mech_items if item != '']
    add_bullet_text(s, Inches(0.7), Inches(1.4), Inches(5.0), Inches(3.8), mech_items, font_size=11, line_spacing=1.3)

    add_shape_bg(s, WHITE, Inches(6.5), Inches(0.8), Inches(5.0), Inches(4.5))
    add_textbox(s, Inches(6.8), Inches(0.9), Inches(4.5), Inches(0.4),
                '差异化优势', font_size=14, color=DARK_TEXT, bold=True)
    diff_items = data.get("differentiation", [
        '根据具体产品填写差异化信息',
        '可通过 --data-file 提供定制数据'
    ])
    add_bullet_text(s, Inches(6.8), Inches(1.4), Inches(4.5), Inches(3.8), diff_items, font_size=11, line_spacing=1.2)

    # Slide 4: 研究设计
    s = new_slide()
    add_section_header(s, '3', f'{data.get("study_code", "关键")} 研究设计')
    design_data = [
        ['项目', '详情'],
        ['研究类型', data.get("study_type", "多中心、随机、对照III期")],
        ['中心数', data.get("centers", "XX家")],
        ['入组时间', data.get("enrollment_period", "X年X月 — X年X月")],
        ['主要入排标准', data.get("eligibility", "标准入排条件")],
        ['随机分组', f'1:1 {data["arm_name"]} | {data["control_arm_name"]}'],
        ['主要终点', data.get("primary_endpoint", "ORR / PFS / OS")],
        ['次要终点', data.get("secondary_endpoints", "次要疗效终点 + 安全性 + PK")],
        ['患者总数', f'{data.get("enrollment", "XX")}例'],
        ['基线特征', data.get("baseline", "标准基线特征")],
    ]
    add_table(s, Inches(0.5), Inches(0.8), Inches(11), Inches(5.2),
              len(design_data), 2, design_data, col_widths=[Inches(3.0), Inches(8.0)])

    # Slide 5: 疗效数据
    s = new_slide()
    add_section_header(s, '4', '疗效数据——显著优势')
    add_kpi_box(s, Inches(0.4), Inches(0.85), Inches(2.7), Inches(0.95),
                kd.get("ORR", "XX%"), f'ORR ({data["arm_name"]})', GREEN_GOOD)
    add_kpi_box(s, Inches(3.3), Inches(0.85), Inches(2.7), Inches(0.95),
                kd.get("control_ORR", "XX%"), f'ORR ({data["control_arm_name"]})', RED_BAD)
    add_kpi_box(s, Inches(6.2), Inches(0.85), Inches(2.7), Inches(0.95),
                kd.get("mPFS", "XXmo"), f'mPFS ({data["arm_name"]})', GREEN_GOOD)
    add_kpi_box(s, Inches(9.1), Inches(0.85), Inches(2.7), Inches(0.95),
                kd.get("control_mPFS", "XXmo"), f'mPFS ({data["control_arm_name"]})', RED_BAD)
    add_kpi_box(s, Inches(0.4), Inches(2.0), Inches(5.6), Inches(0.95),
                f'HR = {kd.get("HR", "0.XX")}', '疾病进展风险降低', ACCENT2)
    add_kpi_box(s, Inches(6.2), Inches(2.0), Inches(5.6), Inches(0.95),
                f'DCR {kd.get("DCR", "XX%")} vs {kd.get("control_DCR", "XX%")}', '疾病控制率', ACCENT)

    add_textbox(s, Inches(0.5), Inches(3.15), Inches(10), Inches(0.4),
                '关键疗效数据', font_size=13, color=DARK_TEXT, bold=True)
    eff_data = [
        ['疗效指标', f'{data["arm_name"]}', f'{data["control_arm_name"]}', '差异'],
        ['ORR', kd.get("ORR", "XX%"), kd.get("control_ORR", "XX%"), kd.get("ORR_diff", "XX%")],
        ['DCR', kd.get("DCR", "XX%"), kd.get("control_DCR", "XX%"), kd.get("DCR_diff", "XX%")],
        ['mPFS', kd.get("mPFS", "XXmo"), kd.get("control_mPFS", "XXmo"), f'HR {kd.get("HR", "0.XX")}'],
        ['mDOR', kd.get("mDOR", "XXmo"), kd.get("control_mDOR", "XXmo"), kd.get("DOR_note", "")],
        ['中位起效时间', kd.get("TTR", "XX"), kd.get("control_TTR", "XX"), kd.get("TTR_note", "")],
        ['OS', kd.get("OS", "随访中"), kd.get("control_OS", "随访中"), kd.get("OS_note", "随访中")],
    ]
    add_table(s, Inches(0.5), Inches(3.6), Inches(11), Inches(3.0),
              len(eff_data), 4, eff_data, col_widths=[Inches(2.5), Inches(3.0), Inches(3.0), Inches(2.5)])

    # Slide 6: 亚组/深度分析
    s = new_slide()
    add_section_header(s, '5', '生存获益与亚组分析')
    add_textbox(s, Inches(0.5), Inches(0.8), Inches(5.5), Inches(0.4),
                '亚组分析——获益一致性', font_size=14, color=DARK_TEXT, bold=True)
    subgroup_items = data.get("subgroup_data", [
        data.get("subgroup_notes", "✅ 各亚组一致获益"),
        "可根据具体数据修改",
        "请使用 --data-file 提供定制内容"
    ])
    add_bullet_text(s, Inches(0.5), Inches(1.3), Inches(11.0), Inches(2.5), subgroup_items, font_size=11)

    add_shape_bg(s, HIGHLIGHT_BG, Inches(0.5), Inches(3.8), Inches(11.0), Inches(1.0))
    highlight_text = data.get("highlight", "")
    if not highlight_text:
        highlight_text = f'💡 {data["product_name"]}在关键疗效指标上显示出临床意义的改善'
    add_textbox(s, Inches(0.7), Inches(3.85), Inches(10.5), Inches(0.8),
                highlight_text, font_size=12, color=DARK_TEXT, bold=True)

    # Slide 7: 安全性
    s = new_slide()
    add_section_header(s, '6', '安全性——可管理的毒性谱')
    add_textbox(s, Inches(0.5), Inches(0.75), Inches(7.0), Inches(0.4),
                '治疗相关不良事件', font_size=12, color=DARK_TEXT, bold=True)

    # Build safety table from data
    safety_data = data.get("safety_table", [
        ['不良事件', f'{data["arm_name"]}', f'{data["control_arm_name"]}'],
        ['任何TRAE', 'XX%', 'XX%'],
        ['Grade ≥3 TRAE', 'XX%', 'XX%'],
        ['严重TRAE', 'XX%', 'XX%'],
        ['TRAE导致停药', 'XX%', 'XX%'],
        ['常见AE 1', 'XX%', 'XX%'],
        ['常见AE 2', 'XX%', 'XX%'],
    ])
    add_table(s, Inches(0.5), Inches(1.15), Inches(5.5), Inches(4.0),
              len(safety_data), 3, safety_data, col_widths=[Inches(2.0), Inches(1.75), Inches(1.75)])

    add_textbox(s, Inches(6.5), Inches(0.75), Inches(5.0), Inches(0.4),
                '安全性关键解读', font_size=14, color=DARK_TEXT, bold=True)
    safety_items = data.get("safety_items", [
        '🔬 主要毒性：',
        '   — 可预测、可管理',
        '',
        '✅ 整体耐受性良好',
        '',
        '⚠️ 关注特殊不良反应',
        '   — 发生率低',
        '',
        '💉 支持治疗体系：',
        '   标准支持治疗方案',
    ])
    add_bullet_text(s, Inches(6.5), Inches(1.2), Inches(5.0), Inches(4.2), safety_items, font_size=10.5, line_spacing=1.2)

    # Slide 8: 受众定制内容
    s = new_slide()
    add_section_header(s, '7', '对临床的实际意义')
    add_textbox(s, Inches(0.5), Inches(0.8), Inches(10), Inches(0.4),
                audience_content["title"], font_size=14, color=DARK_TEXT, bold=True)
    # Split into two columns for readability
    items = audience_content["items"]
    mid = len(items) // 2
    add_bullet_text(s, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.0),
                    items[:mid], font_size=11.5, line_spacing=1.3)
    add_bullet_text(s, Inches(6.5), Inches(1.3), Inches(5.0), Inches(5.0),
                    items[mid:], font_size=11.5, line_spacing=1.3)

    # Slide 9: 竞品对比
    s = new_slide()
    add_section_header(s, '8', '竞品格局与定位')
    comp_data = data.get("comparator_data", DEFAULT_DATA["comparator_data"])
    rows = len(comp_data)
    cols = len(comp_data[0]) if comp_data else 2
    add_table(s, Inches(0.5), Inches(0.8), Inches(11), Inches(3.0),
              rows, cols, comp_data)
    add_textbox(s, Inches(0.5), Inches(4.0), Inches(5.5), Inches(0.4),
                '核心差异化优势', font_size=14, color=DARK_TEXT, bold=True)
    pos_items = data.get("positioning", [
        '🏆 关键临床研究数据亮眼',
        '🎯 针对未满足临床需求',
        '🛡️ 安全性可管理',
        '🔮 进一步研究拓展中'
    ])
    add_bullet_text(s, Inches(0.5), Inches(4.4), Inches(5.5), Inches(2.0), pos_items, font_size=11.5)

    # Slide 10: Take-Home
    s = new_slide()
    add_bg(s, BG_DARK)
    add_shape_bg(s, ACCENT, Emu(0), Emu(0), Emu(80000), prs.slide_height)
    add_textbox(s, Inches(1.5), Inches(0.5), Inches(9), Inches(0.6),
                '⭐ Take-Home Messages', font_size=28, color=WHITE, bold=True)
    add_shape_bg(s, ACCENT2, Inches(1.5), Inches(1.15), Inches(3), Pt(3))

    thm = data.get("take_home", DEFAULT_DATA["take_home"])
    y = Inches(1.5)
    for i, msg in enumerate(thm):
        icon_num = ['❶', '❷', '❸', '❹', '❺'][i] if i < 5 else f'  {i+1}.'
        add_shape_bg(s, ACCENT, Inches(1.8), y, Inches(0.5), Inches(0.5))
        add_textbox(s, Inches(1.85), y + Inches(0.05), Inches(0.4), Inches(0.4),
                    icon_num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, Inches(2.6), y, Inches(9), Inches(0.7),
                    msg, font_size=13, color=RGBColor(0xDD, 0xDD, 0xDD))
        y += Inches(1.1)
    add_shape_bg(s, ACCENT2, Inches(1.5), Inches(6.2), Inches(9.5), Emu(40000))
    add_textbox(s, Inches(2.0), Inches(6.25), Inches(8.5), Inches(0.4),
                f'内部学术交流 | {data["product_name"]} | {data.get("study_code", "关键临床研究")}',
                font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Slide 11: Q&A
    s = new_slide()
    add_bg(s, BG_DARK)
    add_shape_bg(s, ACCENT, Emu(0), Emu(0), prs.slide_width, Emu(70000))
    add_textbox(s, Inches(0.5), Inches(0.1), Inches(10), Inches(0.6),
                'Q&A', font_size=28, color=WHITE, bold=True)
    add_shape_bg(s, ACCENT2, Inches(0.5), Inches(0.75), Inches(2), Pt(3))
    add_textbox(s, Inches(0.5), Inches(1.2), Inches(10), Inches(0.5),
                '感谢聆听！学术讨论欢迎交流', font_size=24, color=WHITE)
    refs = data.get("references", [])
    ref_text = '\n'.join(refs) if refs else \
        f'参考文献：\n请根据实际产品填写相关参考文献\n\n本幻灯片仅供内部学术交流使用'
    add_textbox(s, Inches(0.5), Inches(1.9), Inches(10), Inches(1.5),
                ref_text, font_size=11, color=RGBColor(0xAA, 0xAA, 0xAA))


def load_data_from_file(filepath):
    """从JSON文件加载产品数据"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        print(f"⚠️ 数据文件读取失败: {e}")
        return None


def get_default_data(product_name, indication='xx肿瘤', audience='综合'):
    """生成基于产品名称和适应症的默认数据"""
    data = dict(DEFAULT_DATA)
    data["product_name"] = product_name
    data["indication"] = indication
    
    # 如果用户只输入产品名，给出引导说明
    if product_name != "产品A" and indication == 'xx肿瘤':
        print(f"💡 提示：当前使用模板数据，建议通过 --data-file 提供详细数据")
        print(f"   参考数据模板: python3 generate_dept_meeting_ppt.py --generate-template")
    
    return data


def generate_template_json():
    """生成数据模板JSON"""
    template = {
        "_comment": "科室会PPT数据模板 — 替换占位值即可",
        "product_name": "产品名称",
        "product_code": "产品代号（如BL-XXXX）",
        "drug_class": "药物类别（如双特异性ADC/PD-1抑制剂）",
        "mechanism": "作用机制描述",
        "target": "靶点",
        "antibody_type": "抗体类型（如人源化IgG1）",
        "payload": "载荷（如适用）",
        "study_code": "关键研究代号",
        "indication": "适应症",
        "incidence": "发病率描述",
        "disease": "疾病名称",
        "study_type": "多中心、随机、III期",
        "centers": "XX家",
        "enrollment": "XX例",
        "enrollment_period": "入组时间",
        "eligibility": "主要入排标准描述",
        "primary_endpoint": "PFS/OS/ORR",
        "secondary_endpoints": "次要终点",
        "baseline": "基线特征描述",
        "arm_name": "试验组名称",
        "control_arm_name": "对照组名称",
        "differentiation": [
            "差异化优势1",
            "差异化优势2",
            "差异化优势3"
        ],
        "key_data": {
            "ORR": "XX%",
            "control_ORR": "XX%",
            "ORR_diff": "Δ=XX%",
            "mPFS": "XX个月",
            "control_mPFS": "XX个月",
            "HR": "0.XX",
            "DCR": "XX%",
            "control_DCR": "XX%",
            "DCR_diff": "+XX%",
            "mDOR": "XX个月",
            "control_mDOR": "XX个月",
            "DOR_note": "翻倍",
            "TTR": "1.54个月",
            "control_TTR": "1.54个月",
            "TTR_note": "相同",
            "OS": "随访中",
            "control_OS": "随访中",
            "OS_note": "随访中"
        },
        "subgroup_data": [
            "✅ 亚组1 → 一致获益",
            "✅ 亚组2 → 一致获益"
        ],
        "subgroup_notes": "各亚组一致获益",
        "highlight": "💡 疗效亮点描述",
        "safety_table": [
            ["不良事件", "试验组", "对照组"],
            ["任何TRAE", "XX%", "XX%"],
            ["Grade ≥3 TRAE", "XX%", "XX%"],
            ["严重TRAE", "XX%", "XX%"],
            ["TRAE导致停药", "XX%", "XX%"]
        ],
        "safety_items": [
            "🔬 主要毒性描述",
            "",
            "✅ 整体可管理",
            "",
            "⚠️ 关注特定AE",
            "",
            "💉 支持治疗：标准方案"
        ],
        "comparator_data": [
            ["药物/方案", "类型", "阶段", "ORR", "mPFS"],
            ["本品", "类别", "III期", "XX%", "XXmo"],
            ["对照方案1", "类别", "已有/II期", "XX%", "XXmo"],
            ["对照方案2", "类别", "已有", "XX%", "XXmo"]
        ],
        "positioning": [
            "🏆 优势1",
            "🎯 优势2",
            "🛡️ 优势3",
            "🔮 展望"
        ],
        "take_home": [
            "第一条核心信息",
            "第二条核心信息",
            "第三条核心信息",
            "第四条核心信息",
            "第五条核心信息"
        ],
        "references": [
            "参考文献1",
            "参考文献2"
        ]
    }
    return template


def main():
    parser = argparse.ArgumentParser(description='通用科室会PPT生成器')
    parser.add_argument('--product', default='产品A', help='产品名称')
    parser.add_argument('--indication', default='xx肿瘤', help='适应症')
    parser.add_argument('--audience', default='综合', choices=AUDIENCES, help='目标科室')
    parser.add_argument('--data-file', help='JSON数据文件路径')
    parser.add_argument('--generate-template', action='store_true', help='输出数据模板JSON')

    args = parser.parse_args()

    # 如果需要生成模板
    if args.generate_template:
        tmpl = generate_template_json()
        print(json.dumps(tmpl, ensure_ascii=False, indent=2))
        return

    # 加载数据
    data = None
    if args.data_file:
        data = load_data_from_file(args.data_file)
    if not data:
        data = get_default_data(args.product, args.indication, args.audience)

    # 设置indication（默认数据可能没有）
    if "indication" not in data:
        data["indication"] = args.indication

    init_presentation()
    build_standard_deck(data, audience=args.audience)

    import datetime
    today = datetime.date.today().strftime("%Y%m%d")
    safe_name = data["product_name"].replace("/", "-")
    output = os.path.expanduser(f"~/Desktop/科室会-{safe_name}-{args.audience}-{today}.pptx")
    prs.save(output)
    print(f"✅ PPT已生成：{output}")
    print(f"   总页数：{len(prs.slides)}")
    print(f"   产品：{data['product_name']} | 适应症：{data.get('indication', '')} | 科室：{args.audience}")


if __name__ == "__main__":
    main()
